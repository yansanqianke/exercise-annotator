"""文档解析服务 — 支持 PDF、Word、PPT 文本提取与分块"""

import io

from llama_index.core.node_parser import SentenceSplitter


def extract_text(filename: str, content: bytes) -> str:
    """根据文件扩展名提取文本内容"""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(content)
    elif ext in ("docx", "doc"):
        return _extract_docx(content)
    elif ext in ("pptx", "ppt"):
        return _extract_pptx(content)
    elif ext == "txt":
        return content.decode("utf-8", errors="ignore")
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _extract_pdf(content: bytes) -> str:
    """从 PDF 提取文本"""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def _extract_docx(content: bytes) -> str:
    """从 Word 文档提取文本"""
    from docx import Document

    doc = Document(io.BytesIO(content))
    text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(text_parts)


def _extract_pptx(content: bytes) -> str:
    """从 PPT 提取文本"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)
    return "\n".join(text_parts)


def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> list[str]:
    """按句子边界分块 — 使用 LlamaIndex SentenceSplitter，基于 token 数切分"""
    if not text.strip():
        return []

    splitter = SentenceSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
    )
    return splitter.split_text(text)
